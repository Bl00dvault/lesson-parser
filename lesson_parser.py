import re
import os
from pptx import Presentation

def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    slides = {}

    for slide in presentation.slides:
        slide_num_search = re.search(r"slide(\d+)\.xml$", slide.part.partname)
        slide_num = slide_num_search.group(1)
        slide_data = {"Slide Number": slide_num, "Slide Title": "", "Slide Text": ""}
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Set slide title if not already set
                if not slide_data["Slide Title"] and shape.top <= 300000:
                    slide_data["Slide Title"] = shape.text

                # Append shape text to slide_data["Slide Text"]
                if not shape.text == slide_data["Slide Title"]:
                    slide_data["Slide Text"] += shape.text
                
        # Strip extra line breaks
        slide_data["Slide Text"] = slide_data["Slide Text"].rstrip("\n")

        # Add slide data to slides dictionary
        slides[slide_num] = slide_data

    # Print the dictionary
    # for slide_num, slide_data in slides.items():
    #     print("Slide:", slide_num)
    #     print("Slide Title:", slide_data["Slide Title"])
    #     print("Slide Text:", slide_data["Slide Text"])
    #     print()

    return slides


def write_output_to_file(output_data, input_folder, individual_output_file=None, combined_output_file=None, objectives_only=None):
    if combined_output_file:
        combined_output_path = os.path.join(input_folder, combined_output_file)
        with open(combined_output_path, "w", encoding="utf-8") as combined_file:
            for file_name, slide_data in output_data.items():
                combined_file.write(f"PowerPoint: {file_name}\n")
                for slide_num, slide_content in slide_data.items():
                    combined_file.write(f"Slide: {slide_num}\n")
                    combined_file.write(f"Slide Title: {slide_content['Slide Title']}\n")
                    combined_file.write(f"Slide Text: {slide_content['Slide Text']}\n")
                    combined_file.write("\n")
    
    if objectives_only:
        combined_output_path = os.path.join(input_folder, objectives_only)
        with open(combined_output_path, "w", encoding="utf-8") as combined_file:
            for file_name, slide_data in output_data.items():
                combined_file.write(f"PowerPoint: {file_name}\n")
                for slide_num, slide_content in slide_data.items():
                    # Added slide_num <= 10 to avoid catching recap of objectives at end of lesson
                    if (slide_content["Slide Title"] == "Objectives" or slide_content["Slide Title"] == "Academic Objectives") and int(slide_num) <= 10:
                        combined_file.write(f"Slide: {slide_num}\n")
                        combined_file.write(f"Slide Title: {slide_content['Slide Title']}\n")
                        combined_file.write(f"Slide Text: {slide_content['Slide Text']}\n")
                        combined_file.write("\n")

    elif individual_output_file:
        for file_name, slide_data in output_data.items():
            output_file = os.path.splitext(file_name)[0] + ".txt"
            output_file_path = os.path.join(input_folder, output_file)
            with open(output_file_path, "w", encoding='utf-16') as file:
                # Write the first line with "PowerPoint: " and the filename after the first dash "-"
                filename_parts = file_name.split("-")
                filename_parts = filename_parts[1].split(".")
                if len(filename_parts) > 1:
                    prefix = "PowerPoint: " + filename_parts[0].strip()
                else:
                    prefix = "PowerPoint: " + file_name
                file.write(prefix + "\n")
                for slide_num, slide_content in slide_data.items():
                    if slide_content['Slide Title'] == '':
                        continue
                    elif slide_content['Slide Title'] == 'Print Instructions':
                        continue
                    else:
                        file.write(f"Slide: {slide_num}\n")
                        file.write(f"Slide Title: {slide_content['Slide Title']}\n")
                        file.write(f"Slide Text: {slide_content['Slide Text']}\n")

def extract_text_from_ppt_folder(folder_path):
    slides = {}
    file_list = os.listdir(folder_path)
    for file_name in file_list:
        if file_name.endswith(".pptx") or file_name.endswith(".ppt"):
            file_path = os.path.join(folder_path, file_name)
            slide_data = extract_text_from_ppt(file_path)
            slides[file_name] = slide_data

    return slides

def main():
    current_directory = os.getcwd()
    input_folder = "lessons"
    input_folder_path = os.path.join(current_directory, input_folder)

    output_folder = "text"
    combined_output_file = '_all-lessons.txt'
    objectives_only = '_objectives.txt'
    output_folder_path = os.path.join(current_directory, output_folder)

    all_ppt_data = extract_text_from_ppt_folder(input_folder_path)

    # Set "combined_output_file" to create "all-lessons.txt" which is just the entire output in one text file
    # Set "individual_output_file" to create an individual txt file for each powerpoint
    # Set "objectives_only" to create a "objectives.txt" for just a combined list of acadmic objectives
    write_output_to_file(all_ppt_data, output_folder_path, individual_output_file=None, combined_output_file=None, objectives_only=objectives_only)


if __name__ == '__main__':
    main()

